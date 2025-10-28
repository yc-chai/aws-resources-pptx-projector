from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from collections import defaultdict
import os
import copy
import json

# ==============================
# USER CONFIGURATION (Customize These)
# ==============================
"""
These values can be customized before running the script.
- FILE_INPUT: path to your input JSON file
- FILE_OUTPUT: output PowerPoint file path
- SLIDE_W / SLIDE_H: slide dimensions (in inches)
"""

# --- File / Directory Locations ---
FILE_INPUT = Path("data/sample_aws_resources.json")  # Set your input JSON file path
FILE_OUTPUT = Path(
    "powerpoint/sample_output.pptx"
)  # Set your desired Powerpoint output file path

# --- Content Starting Point ---
START_LEFT = Inches(0.1)  # Put your content to be generated left value
START_TOP = Inches(0.1)  # Put your content to be generated height value

# --- Slide Dimension ---
SLIDE_W = Inches(13.33) - (
    START_LEFT * 2
)  # Set your slide width (default 13.33" for Oracle layout)
SLIDE_H = Inches(5.48)  # Set your silde height

# --- PowerPoint Styling ---
PPTX_SLIDE_LAYOUT = 0  # Set your slide layout (Home > Layout), find the layout type number start with 0 (default 15 for Oracle layout type)
PPTX_FONT_SIZE = Pt(12)  # Set the text font size

# ==============================
# INTERNAL CONFIGURATION ( Do Not Modify Unless Needed )
# ==============================

DIR_AWS_ICON_ITEM = Path("images/aws_icons/")
DIR_AWS_ICON_GROUP = Path("images/aws_group_icons/")

# GROUP ICON
GROUP_ICON_W = Inches(0.35)
GROUP_ICON_H = Inches(0.35)
GROUP_LABEL_TB_W = Inches(2.2)
GROUP_LABEL_TB_H = Inches(0.3)
GROUP_GAP_ICON_LABEL = Inches(0.05)

# GROUP SIZE
GROUP_W = Inches(2)
GROUP_H = Inches(0.5)
GROUP_RATIO_W = 4
GROUP_RATIO_H = 3
GROUP_TARGET_RATIO = GROUP_RATIO_W / GROUP_RATIO_H

# ITEM COMPONENTS
ITEM_ICON_W = Inches(0.5)
ITEM_ICON_H = Inches(0.5)
ITEM_DESC_TB_W = Inches(2)
ITEM_DESC_TB_H = Inches(0.4)
ITEM_GAP_ICON_DESC = Inches(0.05)

# ITEM SIZE
ITEM_W = max(ITEM_ICON_W, ITEM_DESC_TB_W)
ITEM_H = ITEM_ICON_H + ITEM_GAP_ICON_DESC + ITEM_DESC_TB_H

# ITEM & GROUP STYLING
PAD_H = Inches(0.15)  # Padding Horizontal (left-right)
PAD_V = Inches(0.25)  # Padding Vertical (top-bottom)
GAP_H = Inches(0.15)  # Item Gap Horizontal
GAP_V = Inches(0.25)  # Item Gap Vertical

# ITEM ICON STYLING
ITEM_ICON_MAP = {
    "aws_cloud": f"{DIR_AWS_ICON_ITEM}/aws_cloud.png",
    "ebs": f"{DIR_AWS_ICON_ITEM}/ebs.png",
    "ec2": f"{DIR_AWS_ICON_ITEM}/ec2.png",
    "elb": f"{DIR_AWS_ICON_ITEM}/elb.png",
    "igw": f"{DIR_AWS_ICON_ITEM}/igw.png",
    "rds": f"{DIR_AWS_ICON_ITEM}/rds.png",
    "s3": f"{DIR_AWS_ICON_ITEM}/s3.png",
}

# GROUP BORDER STYLING
BORDER_COLOR_MAP = {
    "region": RGBColor(0x49, 0xA1, 0xA5),
    "az": RGBColor(0x49, 0xA1, 0xA5),
    "vpc": RGBColor(0x84, 0x52, 0xF6),
    "private_subnet": RGBColor(0x49, 0xA1, 0xA5),
    "public_subnet": RGBColor(0x82, 0xA0, 0x36),
    "default": RGBColor(0x00, 0x00, 0x00),
}
BORDER_ICON_MAP = {
    "region": f"{DIR_AWS_ICON_GROUP}/region.png",
    "vpc": f"{DIR_AWS_ICON_GROUP}/vpc.png",
    "private_subnet": f"{DIR_AWS_ICON_GROUP}/private_subnet.png",
    "public_subnet": f"{DIR_AWS_ICON_GROUP}/public_subnet.png",
    "default": None,
}
BORDER_DASH_MAP = {
    "region": MSO_LINE_DASH_STYLE.SQUARE_DOT,
    "az": MSO_LINE_DASH_STYLE.DASH,
    "default": MSO_LINE_DASH_STYLE.SOLID,
}

# CALCULATION
SPECIAL_ITEM_CATE = ["igw"]


# ==============================
# DEBUGGING FUNCTION
# ==============================
def print_json(data: dict, indent: int = 2) -> None:
    """Used for debugging, print dict, JSON type console output"""
    print(json.dumps(data, indent=indent))


# ==============================
# TRANSFORM FUNCTION
# ==============================
def load_data(file_path: str) -> list:
    """Open file and return the JSON list data"""
    if not os.path.isfile(file_path):
        exit(f"File not found: {file_path}")

    with open(file_path, "r") as f:
        buffer = f.read()

    try:
        data = json.loads(buffer)
    except ValueError as e:
        exit(f"Invalid JSON in {file_path}: {e}")

    return data


def extract_resources(data: list) -> list:
    """
    Convert AWS JSON into a flat list of collected resources (generic).
    Expected output filelds (if available):
    item, id, name, region, vpc, az, subnet
    """

    def add_result(**item_elements):
        result.append(item_elements)

    result = []
    for region_entry in data:
        region = region_entry.get("region")
        collected = region_entry.get("collected_resources", {})  # Collected_resources

        for item_name, item_data in collected.items():

            # EC2 items
            if item_name == "ec2":
                for ec2_obj in item_data:
                    for resource_type, instances in ec2_obj.items():
                        if resource_type.lower() == "instances":
                            for inst in instances:
                                instance_id = inst.get("InstanceId")
                                vpc_id = inst.get("VpcId")
                                subnet_id = inst.get("SubnetId")
                                az = inst.get("Placement", {}).get("AvailabilityZone")

                                # Default name if no tag
                                name = None
                                tags = inst.get("Tags", [])
                                for tag in tags:
                                    if tag.get("Key") == "Name":
                                        name = tag.get("Value")
                                        break

                                add_result(
                                    item=item_name,
                                    id=instance_id,
                                    name=name,
                                    region=region,
                                    vpc=vpc_id,
                                    az=az,
                                    subnet=subnet_id,
                                )

            # RDS items
            if item_name == "rds":
                for rds_obj in collected["rds"]:
                    rds_id = rds_obj.get("DbiResourceId")  # unique resource ID
                    rds_name = rds_obj.get(
                        "DBInstanceIdentifier"
                    )  # human-readable identifier
                    az = rds_obj.get("AvailabilityZone")

                    db_subnet_group = rds_obj.get("DBSubnetGroup", {})
                    vpc_id = db_subnet_group.get("VpcId")
                    for subnet in db_subnet_group.get("Subnets", []):
                        subnet_az = subnet.get("SubnetAvailabilityZone", {}).get("Name")
                        if subnet_az == az:  # Main AZ
                            add_result(
                                item=item_name,
                                id=rds_id,
                                name=rds_name,
                                region=region,
                                vpc=vpc_id,
                                az=subnet_az,
                                subnet=None,
                            )
                        else:
                            add_result(
                                item=item_name,
                                id=f"{rds_id}-{subnet_az}",
                                name=f"{rds_name}-{subnet_az}",
                                region=region,
                                vpc=vpc_id,
                                az=subnet_az,
                                subnet=None,
                            )

            # ELB item
            if item_name == "loadbalancer":
                for lb in item_data.get("lb_raw", []):
                    lb_id = lb.get("LoadBalancerArn")
                    lb_name = lb.get("LoadBalancerName")
                    vpc_id = lb.get("VpcId")

                    add_result(
                        item="elb",
                        id=lb_id,
                        name=lb_name,
                        region=region,
                        vpc=vpc_id,
                        az=None,
                        subnet=None,
                    )

            if item_name == "network":
                for ig in item_data.get("IGgateway_raw", []):
                    ig_id = ig.get("InternetGatewayId")

                    for attach in ig.get("Attachments", []):
                        if attach.get("State") == "available":
                            vpc_id = attach.get("VpcId")
                            break

                    # name = None
                    # tags = ig.get("Tags", [])
                    # for tag in tags:
                    #     if tag.get("Key") == "Name":
                    #         name = tag.get("Value")
                    #         break
                    # if name is None:
                    #     name = ig_id

                    add_result(
                        item="igw",
                        id=ig_id,
                        region=region,
                        vpc=vpc_id,
                        az=None,
                        subnet=None,
                    )

    return result


def generate_group_items_mapping(items: list) -> list:
    """
    Transform input list of resources into hierarchical group/item mapping.
    - Duplicates group if needed
    - Keeps original IDs
    - Builds cross-linked sharedGroup for VPC & AZ relationship
    """

    def add_node(
        node_id: str,
        node_type: str,
        category: str,
        data: dict = None,
        parent_ids: list = None,
        shared_info: dict = None,
    ) -> None:
        """Add a node (group or item) if not already seen."""
        if not node_id or node_id in seen:
            return

        node = {"id": node_id, "type": node_type, "category": category}
        if data:
            node["data"] = data
        if parent_ids:
            node["parentId"] = parent_ids
        if shared_info:
            node["sharedGroup"] = shared_info

        result.append(node)
        seen.add(node_id)

    def generate_data(**data) -> dict:
        return {k: v for k, v in data.items() if v is not None}

    def generate_shared_info(isPrimaryGroup: bool, groupId: list) -> dict:
        return {"isPrimaryGroup": isPrimaryGroup, "groupId": groupId}

    result = []
    seen = set()

    # Track relationships, collect All VPC-AZ pairings
    vpc_to_azs = defaultdict(set)
    az_to_vpcs = defaultdict(set)

    for item in items:
        vpc = item.get("vpc")
        az = item.get("az")
        if vpc and az:
            vpc_to_azs[vpc].add(az)
            az_to_vpcs[az].add(vpc)

    for item in items:
        region = item.get("region")
        vpc = item.get("vpc")
        az = item.get("az")
        subnet = item.get("subnet")
        iid = item.get("id")
        item_type = item.get("item")
        item_name = item.get("name")

        # Region
        if region:
            add_node(region, "group", "region")

        # VPC & AZ groups with sharedGroup info
        if vpc and az:
            add_node(
                vpc,
                "group",
                "vpc",
                parent_ids=[region] if region else None,
                shared_info=generate_shared_info(True, sorted(list(vpc_to_azs[vpc]))),
            )
            add_node(
                az,
                "group",
                "az",
                parent_ids=[region] if region else None,
                shared_info=generate_shared_info(False, sorted(list(az_to_vpcs[az]))),
            )
        else:
            if vpc:
                add_node(vpc, "group", "vpc", parent_ids=[region] if region else None)
            if az:
                add_node(az, "group", "az", parent_ids=[region] if region else None)

        # Subnet
        if subnet:
            parents = [pid for pid in [vpc, az] if pid]
            if not parents and region:
                parents = [region]
            add_node(subnet, "group", "subnet", parent_ids=parents or None)

        # Items
        parents = []
        if subnet:
            parents.append(subnet)
        else:
            if vpc:
                parents.append(vpc)
            if az:
                parents.append(az)
            if not vpc and not az and region:
                parents.append(region)

        if item_name is not None:
            item_name = item_type + "\n" + item_name

        add_node(
            iid,
            "item",
            item_type or "unknown",
            generate_data(name=item_name),
            parent_ids=parents or None,
        )

    return result


def cal_position_mapping(data: list) -> list:
    """
    Calculate positions for each group and item based on hierarchy.
    Return a list with adding position and style (group & shared-group).
    """

    # Helper functions
    def filter_non_primary_grps() -> list:
        """
        - Filter out non primary groups from data.
        - Modify input data to remove the non primary groups
        - Usually only AvailabilityZones.
        """
        nonlocal data
        result = []
        for n in data:
            shared = n.get("sharedGroup", {})
            if shared and shared.get("isPrimaryGroup") is False:
                result.append(n)
        data = [n for n in data if n not in result]
        return result

    def filter_special_items() -> list:
        """
        - Filter out special items that no needed in the group
        - For example, gateways
        """
        nonlocal data
        result = []
        for n in data:
            n_type = n["type"]
            n_cate = n["category"]
            if n_type == "item" and n_cate in SPECIAL_ITEM_CATE:
                result.append(n)
        data = [n for n in data if n not in result]
        return result

    def find_children(node) -> list:
        """Return child(ren) of a node."""
        nonlocal data
        node_id = node["id"]
        return [n for n in data if node_id in n.get("parentId", [])]

    def find_siblings(node, required_all: bool = False) -> list:
        """Find siblings of a node that have both position and style."""

        def filter_parent_primary_group(node):
            """Return node copy keeping only parentIds whose parents are primary groups."""
            result = dict(node)
            parent_ids = node.get("parentId", [])
            primary_parents = []

            for pid in parent_ids:
                for n in data:
                    if n.get("id") == pid:
                        shared = n.get("sharedGroup", {})
                        if shared.get("isprimaryGroup", True):
                            primary_parents.append(pid)
                        break  # stop one matching parent is found

            result["parentId"] = primary_parents
            return result

        nonlocal data

        # Handle top-level nodes (no parentId)
        parent_ids = node.get("parentId")
        if not parent_ids:
            return [
                n
                for n in data
                if "parentId" not in n
                and n.get("id") != node.get("id")
                and "position" in n
                and "style" in n
            ]

        # Apply primary group filtering if not required_all
        if not required_all:
            parent_ids = filter_parent_primary_group(node).get("parentId")

        parent_set = set(parent_ids)
        siblings = []

        # Iterate through data to find matching siblings
        for n in data:
            if n.get("id") == node.get("id"):
                continue  # skip itself

            n_parent_ids = n.get("parentId")
            if not n_parent_ids:
                continue

            # filter sibling’s parents if not required_all
            if not required_all:
                n_parent_ids = filter_parent_primary_group(n).get("parentId")

            n_parent_set = set(n_parent_ids)

            # only match if parent sets are exactly the same
            match = n_parent_set == parent_set

            # must have both position and style to be valid sibling
            if match and "position" in n and "style" in n:
                siblings.append(n)

        return siblings

    def find_neighbour_siblings(node) -> list:
        nonlocal data
        result = []

        parent = None
        for n in data:
            if n["id"] in node["parentId"]:
                parent = n
                break

        if not parent:
            return result

        if len(parent.get("parentId", [])) == 0:
            return result

        siblings_grp = find_siblings(parent, True)
        node_parent_id = node["parentId"]

        for sibling in siblings_grp:
            sibling_children = find_children(sibling)
            for child in sibling_children:
                # Get child parentIds, including non-primary group
                child_parent_id = child["parentId"]
                if (
                    any(p in child_parent_id for p in node_parent_id)
                    and "position" in child
                ):
                    result.append(child)

        return result

    def sort(nodes: list) -> list:
        """
        Sort a list of nodes within the same hierarchy level based on span and structure rules.

        Sorting logic:
        1. If all nodes have only a "default" span key:
            - Primary sort: ascending by number of parentIds (fewer parents first)
            - Secondary sort: descending by total of "default" span values

        2. Otherwise (mixed span keys present):
            - Step 1: sort by ascending length of parentIds (fewer parents first)
            - Step 2: group nodes by type ('group' before 'item')
            - Step 3: prioritize the first node that has both a non-empty "default" span
                    and additional non-default span keys
            - Step 4: sort remaining nodes by descending total span value (sum of all span lists)
            - Step 5: move nodes that only contain "default" spans to the end
        """

        def total_span(node):
            """Sum of all numbers in all span lists."""
            span = node.get("span", {})
            return sum(sum(v) for v in span.values() if isinstance(v, list))

        def has_only_default_span(node):
            """True if span has only the 'default' key."""
            span = node.get("span", {})
            return set(span.keys()) == {"default"}

        def has_mixed_span(node):
            """True if span has 'default' and at least one other key."""
            span = node.get("span", {})
            keys = list(span.keys())
            return "default" in keys and len(keys) > 1

        def has_nonempty_default(node):
            """True if 'default' span exists and not empty."""
            span = node.get("span", {})
            return bool(span.get("default"))

        def parent_len(node):
            """Return number of parents (0 if no parentId key)."""
            return len(node.get("parentId", []))

        def type_order(node):
            """Lower means higher priority: group < item < others."""
            t = node.get("type", "")
            if t == "group":
                return 0
            if t == "item":
                return 1
            return 2

        if not nodes:
            return nodes

        only_default_level = all(has_only_default_span(n) for n in nodes)

        # All have only default span
        if only_default_level:
            return sorted(nodes, key=lambda n: (parent_len(n), -total_span(n)))

        # Mixed spans
        # Identify the first mixed-span node with non-empty default
        first_default_mixed = next(
            (n for n in nodes if has_mixed_span(n) and has_nonempty_default(n)),
            None,
        )

        # Split into groups
        others = [n for n in nodes if n is not first_default_mixed]
        only_default_nodes = [n for n in others if has_only_default_span(n)]
        non_default_nodes = [n for n in others if n not in only_default_nodes]

        # Sort non-default nodes (excluding the "first mixed") by parent len, type, then total span
        non_default_nodes.sort(
            key=lambda n: (
                parent_len(n),
                type_order(n),
                -total_span(n),
            )
        )

        # Assemble final list
        ordered = []
        if first_default_mixed:
            ordered.append(first_default_mixed)
        ordered.extend(non_default_nodes)
        ordered.extend(only_default_nodes)

        # Finally, ensure global priority 1 and 2 applied to entire list
        # (so parent length and type order still take precedence)
        return sorted(
            ordered,
            key=lambda n: (
                parent_len(n),
                type_order(n),
                (
                    0 if n is first_default_mixed else 1
                ),  # 3rd priority: special node first
            ),
        )

    def get_style(type) -> dict:
        """Return default style depending on type."""
        if type == "item":
            return {"width": ITEM_W, "height": ITEM_H}
        elif type == "group":
            return {"width": GROUP_W, "height": GROUP_H}
        else:
            raise ValueError(f"Type {type} is not in type list")

    def find_parent(node) -> dict:
        parent_id = node.get("parentId")

        if not parent_id:
            return

        nonlocal data

        for n in data:
            if n["id"] in parent_id:
                return n

    def find_smallest_position(data: list) -> dict:
        """
        Return the element with the smallest position (both left and top).
        - 'Smallest' means both `left` and `top` values are strictly less than the current minimum.
        - Assumes every element in `data` has a valid `position` dict with 'left' and 'top' keys.
        """

        # Initialize with the first element
        smallest = data[0]

        for item in data[1:]:
            pos = item["position"]
            smallest_pos = smallest["position"]

            # Update only if left and top are smaller
            if pos["left"] < smallest_pos["left"] or pos["top"] < smallest_pos["top"]:
                smallest = item

        return smallest

    def move(
        offset: int, left: int = None, top: int = None, exception: list = []
    ) -> None:

        if left is None and top is None:
            raise ValueError(f"Must have left or top value.")
        elif left is not None and top is not None:
            raise ValueError(f"Left & Top cannot in same time.")

        nonlocal data

        for n in data:

            if n in exception:
                continue

            if left:
                n_left = n.get("position", {}).get("left", 0)
                n_width = n.get("style", {}).get("width", 0)
                if n_left >= offset:
                    n["position"]["left"] += left
                elif n_left + n_width >= offset:
                    n["style"]["width"] += left
            else:  # top
                n_top = n.get("position", {}).get("top", 0)
                n_height = n.get("style", {}).get("height", 0)
                if n_top >= offset:
                    n["position"]["top"] += top
                elif n_top + n_height >= offset:
                    n["style"]["height"] += top

    def add_style(nodes: list, width: int = None, height: int = None) -> None:

        if width is None and height is None:
            raise ValueError(f"Must have width or height value.")
        elif width is not None and height is not None:
            raise ValueError(f"Width & Height cannot in same time.")

        nonlocal data

        for n in data:

            if n not in nodes:
                continue

            if width:
                n_width = n.get("style", {}).get("width", 0)
                if n_width > 0:
                    n["style"]["width"] += width
            else:  # height
                n_width = n.get("style", {}).get("height", 0)
                if n_width > 0:
                    n["style"]["height"] += height

    def remove_node(data_list: list, removing_nodes: list) -> list:
        result = []

        for n in data_list:
            if n not in removing_nodes:
                result.append(n)
        return result

    def closeness_to_ratio(
        width: float, height: float, target_ratio: float = GROUP_TARGET_RATIO
    ) -> float:
        """Return absolute difference from target ratio (smaller is better)"""
        ratio = width / height
        return abs(ratio - target_ratio)

    def normalize_types(obj):
        if isinstance(obj, dict):
            return {k: normalize_types(v) for k, v in obj.items()}
        elif isinstance(obj, list):
            return [normalize_types(i) for i in obj]
        elif isinstance(obj, (float, int)):
            return int(obj)
        return obj

    def simulate_layout_change(node: dict) -> tuple:
        """
        Reposition children node and generate list of data that calcualted the differences ratio (smaller ratio better layout)
        """
        children = find_children(node)
        children = normalize_types(children)

        if GROUP_TARGET_RATIO >= 1 and len(children) <= 2:
            return False, {}
        elif GROUP_TARGET_RATIO < 1 and len(children) <= 1:
            return False, {}
        children = copy.deepcopy(children)

        # Get most right child node
        most_right_child_node = max(
            children, key=lambda c: c["position"]["left"] + c["style"]["width"]
        )

        # Get it's siblings (all parent must match)
        node_siblings_grp = find_siblings(most_right_child_node, True)
        node_siblings_grp = normalize_types(node_siblings_grp)
        node_siblings_grp = copy.deepcopy(node_siblings_grp)
        if len(node_siblings_grp) == 0:
            return False, {}

        # Get all child position and the node style will be generated
        n_child_left = copy.deepcopy(
            normalize_types(most_right_child_node["position"]["left"])
        )
        n_child_top = copy.deepcopy(
            normalize_types(most_right_child_node["position"]["top"])
        )
        n_child_width = most_right_child_node["style"]["width"]
        n_child_height = most_right_child_node["style"]["height"]

        n_width = n_child_left + n_child_width + PAD_H - node["position"]["left"]
        n_height = (
            max(c["position"]["top"] + c["style"]["height"] for c in children)
            + PAD_V
            - node["position"]["top"]
        )

        # Group siblings by top and sorted
        group_by_top = defaultdict(list)
        s_top = None
        for s in node_siblings_grp:
            s_top = s["position"]["top"]
            group_by_top[s_top].append(s)
        if s_top is not None:
            max_top = max(s["position"]["top"] for s in node_siblings_grp)
            max_height = max(
                s["style"]["height"]
                for s in node_siblings_grp
                if s["position"]["top"] == max_top
            )
            group_by_top[max_top + GAP_V + max_height]  # Adding new row top position
        group_by_top = dict(sorted(group_by_top.items(), key=lambda x: x[0]))

        adding_new_row = False
        child_siblings_top_moving_after_child = 0
        pos_changed = False

        # Calculate every possibility to place the node and make a closer ratio
        keys = list(group_by_top.keys())
        for i, top in enumerate(keys):

            # No needed to calculate initial top
            if top == n_child_top:
                continue

            if i + 1 == len(keys) and pos_changed:
                continue

            top_group = group_by_top[top]
            diff = 0

            # Simulated position
            simulated_child_top = top
            if len(top_group) > 0:  # Top group has siblings
                simulated_child_left = GAP_H + max(
                    s["position"]["left"] + s["style"]["width"] for s in top_group
                )
            else:  # For new row
                simulated_child_left = PAD_H + node["position"]["left"]

            # Get the ending position for the most right child node
            # For easy calculation purpose
            simulated_child_right = simulated_child_left + n_child_width
            simulated_child_bottom = simulated_child_top + n_child_height

            simulated_children = remove_node(children, [most_right_child_node])
            simulated_children = normalize_types(simulated_children)
            simulated_children = copy.deepcopy(simulated_children)

            # Get list of widhts and heights
            widths = [
                c["position"]["left"] + c["style"]["width"] for c in simulated_children
            ] + [simulated_child_right]

            if len(top_group) > 0:
                if simulated_child_top < keys[i + 1] and keys[i + 1] <= (
                    simulated_child_bottom + GAP_V
                ):
                    diff = simulated_child_bottom + GAP_V - keys[i + 1]

                if diff > 0:
                    for c in simulated_children:
                        if simulated_child_top < c["position"]["top"]:
                            c["position"]["top"] += diff
            else:
                for c in simulated_children:
                    if (simulated_child_top - GAP_V) < c["position"]["top"]:
                        c["position"]["top"] += n_child_height + GAP_V
            heights = [
                c["position"]["top"] + c["style"]["height"] for c in simulated_children
            ] + [simulated_child_bottom]

            simulated_n_width = (max(widths) + PAD_H) - node["position"]["left"]
            simulated_n_height = (max(heights) + PAD_V) - node["position"]["top"]

            if closeness_to_ratio(n_width, n_height) > closeness_to_ratio(
                simulated_n_width, simulated_n_height
            ):
                if len(top_group) == 0:
                    adding_new_row = True
                else:
                    pos_changed = True

                n_child_left = simulated_child_left
                n_child_top = simulated_child_top
                n_width = simulated_n_width
                n_height = simulated_n_height
                child_siblings_top_moving_after_child = diff

        if (
            n_child_left == most_right_child_node["position"]["left"]
            and n_child_top == most_right_child_node["position"]["top"]
        ):
            return False, {}
        else:
            return True, {
                "child_nid": most_right_child_node["id"],
                "child_width": n_child_width,
                "child_height": n_child_height,
                "pos_left": n_child_left,
                "pos_top": n_child_top,
                "add_new_row": adding_new_row,
                "child_siblings_offset_top_move": child_siblings_top_moving_after_child,
            }

    def shift_node(node: dict, dx: int, dy: int) -> None:
        """Shift the node and its children."""
        node["position"]["left"] += dx
        node["position"]["top"] += dy

        for c in find_children(node):
            shift_node(c, dx, dy)

    # Main Functions
    def cal_grouping(node) -> int:
        """
        Return integer span units for node
        The returned span is an integer >= 1
        Also, group them in sharedGroups
        """
        # Single item is 1 unit
        if node["type"] == "item":
            return 1

        # Recursively compute child spans
        children = find_children(node)
        spans = {}

        # Detect shared group
        shared_groups = []
        spans = {"default": []}
        if "sharedGroup" in node and node["sharedGroup"]:
            shared_groups = node["sharedGroup"].get("groupId", [])
            for sg in shared_groups:
                spans[sg] = []

        # Process children
        child_span = []
        for child in children:
            child_span = cal_grouping(child)

            # Decide which span bucket to add into
            added = False
            if shared_groups:
                for sg in shared_groups:
                    if sg in child.get("parentId", []):
                        spans[sg].append(child_span)
                        added = True
                        break

            if not added:
                spans["default"].append(child_span)

        # Compute total span (sum of all child spans)
        total_span = sum(sum(v) for v in spans.values())
        node["span"] = spans

        return total_span

    def layout_node(node, left, top, depth: int = 0) -> None:
        """
        Set the layout of the node
        """

        node["position"] = {"left": left, "top": top}
        node["style"] = get_style(node.get("type", None))

        children = find_children(node)
        if not children:  # Always item
            return

        children = sort(children)

        for child in children:

            child_type = child.get("type")

            # For first child -> children[0]
            child_left = left + PAD_H
            child_top = top + PAD_V

            # Siblings are all the same parentId
            siblings_grp = find_siblings(child, True)

            # Siblings in same primary group
            siblings_primary_grp = find_siblings(child)

            # Siblings from parent's neighbour
            siblings_neighbour_grp = find_neighbour_siblings(child)

            # Assign variable for the groups compare
            grp = None

            # horizontal expand first
            if len(siblings_grp) > 0:
                child_top = max(s["position"]["top"] for s in siblings_grp)
                child_left = GAP_H + max(
                    s["position"]["left"] + s["style"]["width"] for s in siblings_grp
                )

            elif len(siblings_neighbour_grp) > 0:
                child_top = min(s["position"]["top"] for s in siblings_neighbour_grp)

            elif len(siblings_primary_grp) > 0:
                child_top = GAP_V + max(
                    s["position"]["top"] + s["style"]["height"]
                    for s in siblings_primary_grp
                )

            layout_node(child, child_left, child_top, depth + 1)

        if depth > 0 and len(children) > 1:
            # simulate_layout_change return dict of the most right node may need
            node_has_better_pos, node_detail = simulate_layout_change(node)

            while node_has_better_pos:
                for c in children:
                    if c["id"] == node_detail["child_nid"]:
                        c["position"]["left"] = node_detail["pos_left"]
                        c["position"]["top"] = node_detail["pos_top"]
                        continue

                    if (
                        node_detail["child_siblings_offset_top_move"] > 0
                        and node_detail["pos_top"] < c["position"]["top"]
                    ):
                        shift_node(c, 0, node_detail["child_siblings_offset_top_move"])

                    if (
                        node_detail["add_new_row"]
                        and c["position"]["top"] >= node_detail["pos_top"]
                    ):
                        shift_node(c, 0, node_detail["child_height"] + GAP_V)

                node_has_better_pos, node_detail = simulate_layout_change(node)

        node["style"]["width"] = (
            max(c["position"]["left"] + c["style"]["width"] for c in children)
            + PAD_H
            - (node["position"]["left"] if len(children) > 0 else 0)
        )
        node["style"]["height"] = (
            max(c["position"]["top"] + c["style"]["height"] for c in children)
            + PAD_V
            - (node["position"]["top"] if len(children) > 0 else 0)
        )

    def layout_non_primary_groups(groups: list) -> None:

        nonlocal data

        # Collection position, preparing for move purpose
        left_collections = defaultdict(list)
        top_collections = defaultdict(list)

        for grp in groups:
            children = find_children(grp)

            if len(children) == 0:
                continue

            smallest_pos_node = find_smallest_position(children)
            smallest_pos_node_parent = find_parent(smallest_pos_node)

            if not smallest_pos_node_parent:
                continue

            pos_top = smallest_pos_node["position"]["top"]
            pos_left = smallest_pos_node_parent["position"]["left"]

            grp_height = (
                max(c["position"]["top"] + c["style"]["height"] for c in children)
                + PAD_V / 2
                - pos_top
            )
            grp_width = (
                max(c["position"]["left"] + c["style"]["width"] for c in children)
                + PAD_H / 2
                - pos_left
            )

            grp["position"] = {"top": pos_top, "left": pos_left}
            grp["style"] = {"width": grp_width, "height": grp_height}

            left_collections[pos_left].append(grp)
            top_collections[pos_top].append(grp)

            data.append(grp)

        left_collections = dict(sorted(left_collections.items()))
        top_collections = dict(sorted(top_collections.items()))

        for pos, grp_nodes in left_collections.items():
            move(pos, left=PAD_H / 2, exception=grp_nodes)
            add_style(grp_nodes, width=PAD_H / 2)

            for node in grp_nodes:
                n_width = node["style"]["width"]
                move(pos + n_width, left=PAD_H / 2)

        for pos, grp_nodes in top_collections.items():
            move(pos, top=PAD_V / 2, exception=grp_nodes)
            add_style(grp_nodes, height=PAD_V / 3)

            for node in grp_nodes:
                n_height = node["style"]["height"]
                move(pos + n_height, top=PAD_V / 2)

    def layout_special_items(items: list) -> None:

        nonlocal data

        parent_collection = defaultdict(list)

        for item in items:
            parent_ids = item.get("parentId", [])
            for pid in parent_ids:
                parent_collection[pid].append(item)

        parent_keys = list(parent_collection.keys())
        for n in data:

            n_id = n["id"]
            if n_id not in parent_keys:
                continue

            # Get item left, top & width
            n_left = n["position"]["left"]
            n_top = n["position"]["top"]
            n_width = n["style"]["width"]

            # Get item list under this parent node
            s_items = parent_collection[n_id]

            # For positioning
            n_start = n_left
            n_end = n_left + n_width
            s_items_len = len(s_items)

            i = 1
            for s_it in s_items:
                pos_left = (
                    n_start + ((n_end - n_start) / (s_items_len + 1) * i) - (ITEM_W / 2)
                )
                pos_top = n_top - (ITEM_ICON_H / 2)

                s_it["position"] = {"left": pos_left, "top": pos_top}
                s_it["style"] = {"width": ITEM_W, "height": ITEM_H}

                data.append(s_it)
                i += 1

    # Filter out non primary groups and special items
    non_primary_grps = filter_non_primary_grps()
    special_items = filter_special_items()

    # Find root item
    root_nodes = [n for n in data if not n.get("parentId")]

    for root in root_nodes:
        cal_grouping(root)

    root_nodes = sort(root_nodes)

    current_top = START_TOP
    current_left = START_LEFT

    for root in root_nodes:
        layout_node(node=root, left=current_left, top=current_top)
        current_left += root["style"]["width"] + GAP_H

    layout_non_primary_groups(non_primary_grps)
    layout_special_items(special_items)

    return data


def generate_pptx(data: list) -> None:
    """
    Generate Powerpoint shapes from data then save
    """
    if FILE_OUTPUT.exists():
        prs = Presentation(pptx=FILE_OUTPUT)
    else:
        prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[PPTX_SLIDE_LAYOUT])

    def add_border_box(category: str, left, top, width, height, text: str = None):
        """
        Add a bordered box of given type at given position.
        """

        def get_color(type: str = category) -> RGBColor:
            """Return RGBColor based on category type"""
            return BORDER_COLOR_MAP.get(type.lower(), BORDER_COLOR_MAP["default"])

        def get_icon(type: str = category) -> str:
            """Return group icon (If available)"""
            return BORDER_ICON_MAP.get(type.lower(), BORDER_ICON_MAP["default"])

        def get_dash(type: str = category) -> MSO_LINE_DASH_STYLE:
            """Return dash style based on category type"""
            return BORDER_DASH_MAP.get(type.lower(), BORDER_DASH_MAP["default"])

        # Add border shape
        grp_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        grp_shape.fill.background()  # Remove fill (transparent)
        grp_shape.line.width = Pt(1.25)  # Border width: 1.25 pt

        # Setting border the shape
        grp_shape.line.color.rgb = get_color()
        grp_shape.line.dash_style = get_dash()

        # Group Icon (image) & Label (textbox)
        grp_icon = get_icon()

        grp_label_left = left
        grp_label_top = top

        if grp_icon is not None:
            grp_label_left += GROUP_ICON_W + GROUP_GAP_ICON_LABEL

            grp_icon_left = left
            grp_icon_top = top
            slide.shapes.add_picture(
                grp_icon, grp_icon_left, grp_icon_top, GROUP_ICON_W, GROUP_ICON_H
            )

        if text is not None:
            grp_label = slide.shapes.add_textbox(
                grp_label_left, grp_label_top, GROUP_LABEL_TB_W, GROUP_LABEL_TB_H
            )

            label_frame = grp_label.text_frame
            label_frame.text = text
            label_frame.word_wrap = True
            label_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            for paragraph in label_frame.paragraphs:
                paragraph.font.size = PPTX_FONT_SIZE

    def add_item_box(category: str, left, top, text: str):
        """Add image + text inside a framed box at given position"""

        def get_icon(type: str = category) -> str:
            """Return item icon"""
            result = ITEM_ICON_MAP.get(type.lower(), None)
            if result:
                return result
            else:
                raise ValueError(f"Icon type: {type} is not found")

        # Image Position
        img_left = left + (ITEM_W / 2 - ITEM_ICON_W / 2)
        img_top = top
        img_path = get_icon()

        slide.shapes.add_picture(
            img_path, left=img_left, top=img_top, width=ITEM_ICON_W, height=ITEM_ICON_H
        )

        if text is None:
            return

        # Textbox Position
        tb_left = left
        tb_top = top + ITEM_ICON_H + ITEM_GAP_ICON_DESC

        textbox = slide.shapes.add_textbox(
            left=tb_left, top=tb_top, width=ITEM_DESC_TB_W, height=ITEM_DESC_TB_H
        )
        frame = textbox.text_frame
        frame.text = text
        frame.word_wrap = True
        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        for paragraph in frame.paragraphs:
            paragraph.font.size = PPTX_FONT_SIZE
            paragraph.alignment = PP_ALIGN.CENTER

    def save_file() -> None:
        """Save file with checking the correct file path"""
        file_path = FILE_OUTPUT
        try:
            file_path.parent.mkdir(
                parents=True, exist_ok=True
            )  # Ensure directory exists
            prs.save(file_path)
            print(f"Presentation saved to {file_path}")
        except Exception as e:
            print(f"Error saving presentation: {e}")

    for node in data:
        node_id = node["id"]
        node_type = node["type"]
        node_cate = node["category"]

        # Position
        pos_left = node["position"]["left"]
        pos_top = node["position"]["top"]

        if node_type == "group":
            # Style (Width & Height)
            style_w = node["style"]["width"]
            style_h = node["style"]["height"]

            add_border_box(node_cate, pos_left, pos_top, style_w, style_h, node_id)
        else:  # node_type == "item"
            node_data = node.get("data", {})

            data_name = node_data.get("name", None)
            # if data_name is None:
            #     data_name = node_id

            add_item_box(node_cate, pos_left, pos_top, data_name)
    save_file()


def main() -> None:
    input_file = FILE_INPUT

    json_data = load_data(input_file)
    # print_json(json_data)

    flat_data = extract_resources(json_data)
    # print_json(flat_data)

    grouped_items = generate_group_items_mapping(flat_data)
    # print_json(grouped_items)

    positioned_items = cal_position_mapping(grouped_items)
    # print_json(positioned_items)

    generate_pptx(positioned_items)


if __name__ == "__main__":
    main()
